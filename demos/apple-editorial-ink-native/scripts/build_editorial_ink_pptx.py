"""生成 Apple 历年财报解读的 Editorial Ink 风格原生 PPTX。

本脚本服务于 `ppt-polished-deck-collab` 的风格迁移测试：读取 workspace 中
已经处理好的 Apple FY2021-FY2025 财务 CSV 与派生 `slide_specs.yaml`，
把 guizang-ppt-skill Style A 的“电子杂志 × 电子墨水”视觉语言翻译为
PowerPoint 原生文本框、形状、线条、表格和图形。整体流程为：
1. 载入财务、产品、地区、毛利率和来源数据；
2. 计算页面所需的同比、占比、利润率和资本回报指标；
3. 使用 native PPT shapes 绘制杂志式 hero、KPI 卡、shape chart、流程图和表格；
4. 将可编辑 PPTX 写入 `build/pptx/` 与 `final/`。
"""

from __future__ import annotations

import json
import shutil
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

import pandas as pd
import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


WORKSPACE = Path(__file__).resolve().parents[1]
DATA_DIR = WORKSPACE / "data" / "processed"
BUILD_DIR = WORKSPACE / "build"
PPTX_DIR = BUILD_DIR / "pptx"
FINAL_DIR = WORKSPACE / "final"
SPEC_PATH = BUILD_DIR / "generated" / "slide_specs.yaml"

OUT_PPTX = PPTX_DIR / "apple_editorial_ink_native_test.pptx"
FINAL_PPTX = FINAL_DIR / "apple_editorial_ink_native_test.pptx"
SUMMARY_PATH = BUILD_DIR / "generated" / "build_summary.json"

PAGE_W = 13.333
PAGE_H = 7.5
SOURCE_TEXT = "数据来源：SEC companyfacts、Apple FY2021-FY2025 10-K；金额单位除特别说明外为十亿美元。"
DISCLAIMER = "仅供学术交流与 skill 测试使用，不构成投资建议。"

INK = RGBColor(10, 31, 61)
INK_TINT = RGBColor(21, 42, 74)
INK_SOFT = RGBColor(41, 62, 91)
PAPER = RGBColor(241, 243, 245)
PAPER_TINT = RGBColor(228, 232, 236)
PAPER_WARM = RGBColor(247, 244, 238)
WHITE = RGBColor(252, 252, 249)
COPPER = RGBColor(188, 132, 79)
DARK_COPPER = RGBColor(131, 88, 53)
CORAL = RGBColor(174, 47, 68)
CELADON = RGBColor(105, 146, 135)
MIST = RGBColor(188, 198, 207)
MUTED = RGBColor(98, 111, 130)
SILVER = RGBColor(210, 216, 222)
TABLE_ALT = RGBColor(236, 239, 242)

FONT_SERIF_CN = "Songti SC"
FONT_SERIF_EN = "Times New Roman"
FONT_SANS = "Helvetica Neue"
FONT_MONO = "Menlo"


@dataclass(frozen=True)
class TextStyle:
    """定义一组可复用的文本风格参数。"""

    size: float
    color: RGBColor
    family: str = "sans"
    bold: bool = False
    italic: bool = False
    align: PP_ALIGN = PP_ALIGN.LEFT
    line_spacing: float = 1.0
    space_before: float = 0.0
    space_after: float = 0.0


def load_inputs() -> dict[str, pd.DataFrame | dict]:
    """读取 slide specs 与处理后的财务数据。"""

    return {
        "spec": yaml.safe_load(SPEC_PATH.read_text(encoding="utf-8")),
        "financials": pd.read_csv(DATA_DIR / "apple_financials_fy2021_fy2025.csv"),
        "products": pd.read_csv(DATA_DIR / "apple_product_net_sales_fy2021_fy2025.csv"),
        "regions": pd.read_csv(DATA_DIR / "apple_region_net_sales_fy2021_fy2025.csv"),
        "gross_margin": pd.read_csv(DATA_DIR / "apple_gross_margin_by_type_fy2021_fy2025.csv"),
        "sources": pd.read_csv(DATA_DIR / "apple_10k_sources.csv"),
    }


def usd_bn(value: float) -> float:
    """把美元金额转换为十亿美元。"""

    return float(value) / 1_000_000_000


def usd_mn_to_bn(value: float) -> float:
    """把百万美元转换为十亿美元。"""

    return float(value) / 1_000


def pct(value: float) -> float:
    """把比例转换为百分数。"""

    return float(value) * 100


def yoy(current: float, previous: float) -> float:
    """计算同比变化率。"""

    return float(current) / float(previous) - 1


def color_tuple(color: RGBColor) -> tuple[int, int, int]:
    """把 RGBColor 转为三元组。"""

    return color[0], color[1], color[2]


def blend(fg: RGBColor, bg: RGBColor, alpha: float) -> RGBColor:
    """按 alpha 把前景色混合到背景色。"""

    return RGBColor(
        int(round(fg[0] * alpha + bg[0] * (1 - alpha))),
        int(round(fg[1] * alpha + bg[1] * (1 - alpha))),
        int(round(fg[2] * alpha + bg[2] * (1 - alpha))),
    )


def make_prs() -> Presentation:
    """创建宽屏 Presentation。"""

    prs = Presentation()
    prs.slide_width = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    return prs


def blank_slide(prs: Presentation):
    """添加空白页。"""

    return prs.slides.add_slide(prs.slide_layouts[6])


def set_transparency(obj, value: float) -> None:
    """设置 fill 或 line 的透明度，数值范围为 0 到 100。"""

    obj.transparency = value


def font_names(family: str) -> tuple[str, str]:
    """按风格族返回 latin 与 east-asian 字体名。"""

    if family == "serif":
        return FONT_SERIF_EN, FONT_SERIF_CN
    if family == "mono":
        return FONT_MONO, FONT_MONO
    return FONT_SANS, FONT_SANS


def apply_run_font(run, style: TextStyle) -> None:
    """为 run 写入字号、颜色、字重和中英文字体槽位。"""

    latin, east_asia = font_names(style.family)
    run.font.name = latin
    run.font.size = Pt(style.size)
    run.font.bold = style.bold
    run.font.italic = style.italic
    run.font.color.rgb = style.color
    rpr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", latin), ("a:ea", east_asia), ("a:cs", latin)):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", typeface)


def apply_paragraph_style(paragraph, style: TextStyle) -> None:
    """为 paragraph 应用基础排版规则。"""

    paragraph.alignment = style.align
    paragraph.line_spacing = style.line_spacing
    paragraph.space_before = Pt(style.space_before)
    paragraph.space_after = Pt(style.space_after)
    for run in paragraph.runs:
        apply_run_font(run, style)


def set_textbox_margins(shape, margin: float = 0.03, anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    """设置文本框边距与垂直锚点。"""

    tf = shape.text_frame
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = anchor
    tf.word_wrap = True


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    style: TextStyle,
    *,
    margin: float = 0.03,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    rotation: float | None = None,
):
    """添加单段文本框。"""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_textbox_margins(box, margin=margin, anchor=anchor)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = line or fill
    elif line is not None:
        box.line.color.rgb = line
    if rotation is not None:
        box.rotation = rotation
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    apply_paragraph_style(paragraph, style)
    return box


def add_multiline_text(
    slide,
    lines: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    style: TextStyle,
    *,
    bullet: bool = False,
    margin: float = 0.04,
):
    """添加多段文本。"""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_textbox_margins(box, margin=margin)
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = f"· {line}" if bullet else line
        apply_paragraph_style(paragraph, style)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    line: RGBColor | None = None,
    transparency: float = 0,
    line_width: float = 0.6,
):
    """添加原生矩形。"""

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    set_transparency(shape.fill, transparency)
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: float = 0.6):
    """添加一条原生直线。"""

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_oval(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    line: RGBColor,
    *,
    fill: RGBColor | None = None,
    transparency: float = 100,
    line_width: float = 0.8,
):
    """添加椭圆或环线，用于模拟电子墨水流体背景。"""

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        set_transparency(shape.fill, transparency)
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def set_background(slide, fill: RGBColor) -> None:
    """设置页面背景色。"""

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = fill


def draw_fluid_background(slide, theme: str) -> None:
    """绘制原生形状版的电子墨水背景。"""

    if theme == "dark":
        line_color = blend(PAPER, INK, 0.22)
        accent = blend(COPPER, INK, 0.55)
    else:
        line_color = blend(INK, PAPER, 0.13)
        accent = blend(COPPER, PAPER, 0.42)

    ovals = [
        (-1.5, -1.0, 6.6, 6.2, 1.0),
        (7.6, -1.4, 7.1, 6.4, 0.8),
        (4.4, 3.3, 8.8, 5.8, 0.7),
        (-2.0, 4.5, 5.4, 3.8, 0.55),
    ]
    for x, y, w, h, width in ovals:
        add_oval(slide, x, y, w, h, line_color, line_width=width)
    for offset in (0.0, 0.18, 0.38):
        add_oval(slide, 8.55 + offset, 0.35 + offset * 0.6, 3.2, 2.6, accent, line_width=0.55)


def draw_dot_field(slide, theme: str, x: float, y: float, cols: int, rows: int, step: float = 0.16) -> None:
    """绘制点阵，模拟 HTML 模板中的像素化节奏。"""

    color = blend(PAPER, INK, 0.28) if theme == "dark" else blend(INK, PAPER, 0.16)
    for row in range(rows):
        for col in range(cols):
            if (row + col) % 3 == 0:
                add_rect(slide, x + col * step, y + row * step, 0.015, 0.015, color, transparency=18, line_width=0)


def add_chrome(slide, theme: str, page_no: int, left: str, right: str | None = None) -> None:
    """添加 guizang 风格的页眉与页脚元数据。"""

    color = PAPER_TINT if theme == "dark" else INK_SOFT
    muted = blend(color, INK if theme == "dark" else PAPER, 0.58)
    right_text = right or f"VOL. TEST · {page_no:02d} / 10"
    meta_style = TextStyle(6.8, muted, "mono", align=PP_ALIGN.LEFT)
    add_text(slide, left.upper(), 0.56, 0.32, 5.4, 0.18, meta_style, margin=0)
    add_text(slide, right_text.upper(), 9.55, 0.32, 3.2, 0.18, TextStyle(6.8, muted, "mono", align=PP_ALIGN.RIGHT), margin=0)
    add_line(slide, 0.56, 6.96, 12.78, 6.96, muted, 0.35)
    add_text(slide, SOURCE_TEXT, 0.56, 7.06, 7.8, 0.20, TextStyle(5.8, muted, "sans"), margin=0)
    add_text(slide, DISCLAIMER, 9.10, 7.06, 3.68, 0.20, TextStyle(5.8, muted, "sans", align=PP_ALIGN.RIGHT), margin=0)


def add_page_title(slide, theme: str, page_no: int, kicker: str, title: str, subtitle: str | None = None) -> None:
    """添加正文页标题区。"""

    color = PAPER if theme == "dark" else INK
    muted = blend(color, INK if theme == "dark" else PAPER, 0.62)
    add_chrome(slide, theme, page_no, kicker)
    add_text(slide, kicker, 0.62, 0.78, 4.8, 0.20, TextStyle(7.6, muted, "mono"), margin=0)
    add_text(
        slide,
        title,
        0.56,
        1.02,
        10.9,
        0.80,
        TextStyle(18.0, color, "serif", bold=True, line_spacing=1.0),
        margin=0,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.60,
            1.82,
            9.7,
            0.38,
            TextStyle(9.4, muted, "sans", line_spacing=1.15),
            margin=0,
        )
    add_text(slide, f"{page_no:02d}", 11.95, 0.86, 0.72, 0.36, TextStyle(16, muted, "mono", align=PP_ALIGN.RIGHT), margin=0)


def add_kpi_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    value: str,
    note: str,
    theme: str,
    accent: RGBColor = COPPER,
) -> None:
    """添加杂志式 KPI 卡片。"""

    if theme == "dark":
        fill = blend(PAPER, INK, 0.08)
        line = blend(PAPER, INK, 0.22)
        text = PAPER
        muted = blend(PAPER, INK, 0.55)
    else:
        fill = blend(PAPER, WHITE, 0.86)
        line = blend(INK, PAPER, 0.16)
        text = INK
        muted = MUTED
    add_rect(slide, x, y, w, h, fill, line=line, transparency=5 if theme == "dark" else 0, line_width=0.6)
    add_rect(slide, x, y, 0.045, h, accent, line=accent, line_width=0)
    add_text(slide, label.upper(), x + 0.15, y + 0.12, w - 0.24, 0.16, TextStyle(6.4, muted, "mono"), margin=0)
    add_text(slide, value, x + 0.14, y + 0.35, w - 0.26, 0.36, TextStyle(17.2, text, "serif", bold=True), margin=0)
    add_text(slide, note, x + 0.15, y + 0.80, w - 0.24, 0.26, TextStyle(7.2, muted, "sans", line_spacing=1.15), margin=0)


def add_statement_block(slide, x: float, y: float, w: float, h: float, title: str, body: str, theme: str) -> None:
    """添加一块带左边线的观点文本。"""

    color = PAPER if theme == "dark" else INK
    muted = blend(color, INK if theme == "dark" else PAPER, 0.62)
    add_rect(slide, x, y, 0.035, h, COPPER, line=COPPER, line_width=0)
    add_text(slide, title, x + 0.16, y + 0.02, w - 0.20, 0.28, TextStyle(11.2, color, "serif", bold=True), margin=0)
    add_text(slide, body, x + 0.16, y + 0.39, w - 0.20, h - 0.35, TextStyle(8.4, muted, "sans", line_spacing=1.28), margin=0)


def add_year_axis(slide, x: float, y: float, w: float, years: list[int], theme: str) -> None:
    """添加五年时间轴。"""

    color = PAPER if theme == "dark" else INK_SOFT
    muted = blend(color, INK if theme == "dark" else PAPER, 0.62)
    add_line(slide, x, y, x + w, y, muted, 0.55)
    for idx, year in enumerate(years):
        px = x + idx * w / (len(years) - 1)
        add_rect(slide, px - 0.018, y - 0.018, 0.036, 0.036, COPPER if year == 2025 else muted, line_width=0)
        add_text(slide, str(year), px - 0.22, y + 0.06, 0.44, 0.15, TextStyle(6.2, muted, "mono", align=PP_ALIGN.CENTER), margin=0)


def chart_coordinates(values: list[float], x: float, y: float, w: float, h: float, *, y_min: float | None = None, y_max: float | None = None) -> list[tuple[float, float]]:
    """把数值映射到图表区域坐标。"""

    min_v = min(values) if y_min is None else y_min
    max_v = max(values) if y_max is None else y_max
    span = max(max_v - min_v, 1e-9)
    coords: list[tuple[float, float]] = []
    for idx, value in enumerate(values):
        px = x + idx * w / (len(values) - 1 if len(values) > 1 else 1)
        py = y + h - (value - min_v) / span * h
        coords.append((px, py))
    return coords


def add_native_line_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    years: list[int],
    series: list[tuple[str, list[float], RGBColor]],
    theme: str,
    *,
    y_min: float | None = None,
    y_max: float | None = None,
    value_suffix: str = "",
) -> None:
    """用原生线条和节点绘制可编辑趋势图。"""

    bg = blend(PAPER, WHITE, 0.55) if theme != "dark" else blend(PAPER, INK, 0.06)
    line_color = blend(INK, PAPER, 0.14) if theme != "dark" else blend(PAPER, INK, 0.22)
    text = INK if theme != "dark" else PAPER
    muted = MUTED if theme != "dark" else blend(PAPER, INK, 0.55)
    add_rect(slide, x, y, w, h, bg, line=line_color, transparency=0 if theme != "dark" else 12)
    plot_x, plot_y, plot_w, plot_h = x + 0.36, y + 0.38, w - 0.66, h - 0.82
    for idx in range(4):
        gy = plot_y + idx * plot_h / 3
        add_line(slide, plot_x, gy, plot_x + plot_w, gy, line_color, 0.25)
    all_values = [value for _, values, _ in series for value in values]
    min_v = min(all_values) if y_min is None else y_min
    max_v = max(all_values) if y_max is None else y_max
    for name, values, color in series:
        coords = chart_coordinates(values, plot_x, plot_y, plot_w, plot_h, y_min=min_v, y_max=max_v)
        for (x1, y1), (x2, y2) in zip(coords, coords[1:]):
            add_line(slide, x1, y1, x2, y2, color, 1.25)
        for idx, (px, py) in enumerate(coords):
            add_oval(slide, px - 0.045, py - 0.045, 0.09, 0.09, color, fill=color, transparency=0, line_width=0)
            if idx == len(coords) - 1:
                add_text(
                    slide,
                    f"{values[idx]:.1f}{value_suffix}",
                    px - 0.10,
                    py - 0.30,
                    0.76,
                    0.16,
                    TextStyle(6.4, text, "mono"),
                    margin=0,
                )
        legend_y = y + h - 0.26
        legend_x = x + 0.36 + series.index((name, values, color)) * 1.65
        add_rect(slide, legend_x, legend_y + 0.05, 0.12, 0.035, color, line=color, line_width=0)
        add_text(slide, name, legend_x + 0.18, legend_y, 1.25, 0.14, TextStyle(6.2, muted, "mono"), margin=0)
    for idx, year in enumerate(years):
        px = plot_x + idx * plot_w / (len(years) - 1)
        add_text(slide, str(year)[2:], px - 0.12, y + h - 0.50, 0.24, 0.12, TextStyle(5.6, muted, "mono", align=PP_ALIGN.CENTER), margin=0)


def add_native_bar_chart(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    years: list[int],
    values: list[float],
    theme: str,
    *,
    color: RGBColor = COPPER,
    label: str = "",
    value_format: str = "{:.0f}",
) -> None:
    """用原生矩形绘制可编辑柱形图。"""

    max_v = max(values)
    muted = blend(PAPER, INK, 0.55) if theme == "dark" else MUTED
    baseline = y + h - 0.40
    chart_h = h - 0.78
    add_line(slide, x, baseline, x + w, baseline, blend(INK, PAPER, 0.18) if theme != "dark" else blend(PAPER, INK, 0.24), 0.45)
    bar_w = min(0.34, w / len(values) * 0.42)
    for idx, value in enumerate(values):
        px = x + idx * w / len(values) + w / len(values) * 0.5 - bar_w / 2
        bh = value / max_v * chart_h
        fill = color if idx == len(values) - 1 else blend(color, PAPER if theme != "dark" else INK, 0.55)
        add_rect(slide, px, baseline - bh, bar_w, bh, fill, line=fill, line_width=0)
        add_text(slide, value_format.format(value), px - 0.13, baseline - bh - 0.18, bar_w + 0.26, 0.12, TextStyle(5.6, muted, "mono", align=PP_ALIGN.CENTER), margin=0)
        add_text(slide, str(years[idx])[2:], px - 0.06, baseline + 0.08, bar_w + 0.12, 0.12, TextStyle(5.6, muted, "mono", align=PP_ALIGN.CENTER), margin=0)
    if label:
        add_text(slide, label.upper(), x, y, w, 0.16, TextStyle(6.0, muted, "mono"), margin=0)


def add_stacked_strip(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    parts: list[tuple[str, float, RGBColor]],
    theme: str,
) -> None:
    """绘制产品收入构成的横向堆叠条。"""

    total = sum(value for _, value, _ in parts)
    current_x = x
    text = PAPER if theme == "dark" else INK
    muted = blend(PAPER, INK, 0.55) if theme == "dark" else MUTED
    for name, value, color in parts:
        segment_w = w * value / total
        add_rect(slide, current_x, y, segment_w, h, color, line=color, line_width=0)
        if segment_w > 1.05:
            add_text(slide, name, current_x + 0.10, y + 0.10, segment_w - 0.18, 0.18, TextStyle(7.0, WHITE, "mono"), margin=0)
            add_text(slide, f"{value:.1f}B", current_x + 0.10, y + 0.34, segment_w - 0.18, 0.20, TextStyle(11.2, WHITE, "serif", bold=True), margin=0)
        current_x += segment_w
    add_text(slide, "FY2025 PRODUCT / SERVICE MIX", x, y - 0.28, 3.6, 0.16, TextStyle(6.4, muted, "mono"), margin=0)
    add_text(slide, f"Total {total:.1f}B", x + w - 1.6, y - 0.28, 1.6, 0.16, TextStyle(6.4, muted, "mono", align=PP_ALIGN.RIGHT), margin=0)
    add_line(slide, x, y + h + 0.16, x + w, y + h + 0.16, muted, 0.35)
    add_text(slide, "构成条为原生 PowerPoint shape，可直接选中编辑颜色、宽度和标签。", x, y + h + 0.24, w, 0.16, TextStyle(6.2, muted, "sans"), margin=0)
    add_text(slide, "Services 增长带来更高质量的收入，但 iPhone 仍承担超过一半收入。", x, y + h + 0.52, w, 0.22, TextStyle(8.4, text, "serif", bold=True), margin=0)


def add_horizontal_bar_matrix(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    rows: list[tuple[str, float, float]],
    theme: str,
) -> None:
    """绘制地区收入排序条与同比矩阵。"""

    max_value = max(value for _, value, _ in rows)
    text = INK if theme != "dark" else PAPER
    muted = MUTED if theme != "dark" else blend(PAPER, INK, 0.55)
    bar_x = x + 1.65
    bar_w = w - 3.55
    row_h = h / len(rows)
    add_text(slide, "REGION", x, y - 0.20, 1.1, 0.14, TextStyle(6.2, muted, "mono"), margin=0)
    add_text(slide, "FY2025 NET SALES", bar_x, y - 0.20, 1.7, 0.14, TextStyle(6.2, muted, "mono"), margin=0)
    add_text(slide, "YOY", x + w - 0.75, y - 0.20, 0.75, 0.14, TextStyle(6.2, muted, "mono", align=PP_ALIGN.RIGHT), margin=0)
    for idx, (name, value, growth) in enumerate(rows):
        cy = y + idx * row_h
        add_line(slide, x, cy + row_h - 0.05, x + w, cy + row_h - 0.05, blend(INK, PAPER, 0.12), 0.25)
        add_text(slide, name, x, cy + 0.08, 1.50, 0.16, TextStyle(7.0, text, "sans", bold=name == "Greater China"), margin=0)
        color = CORAL if growth < 0 else (COPPER if idx < 2 else CELADON)
        fill_w = bar_w * value / max_value
        add_rect(slide, bar_x, cy + 0.10, fill_w, 0.20, color, line=color, line_width=0)
        add_text(slide, f"{value:.1f}B", x + w - 1.66, cy + 0.06, 0.76, 0.15, TextStyle(6.4, muted, "mono", align=PP_ALIGN.RIGHT), margin=0)
        add_text(slide, f"{growth:+.1f}%", x + w - 0.80, cy + 0.06, 0.80, 0.15, TextStyle(6.4, color, "mono", align=PP_ALIGN.RIGHT, bold=True), margin=0)


def add_native_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    headers: list[str],
    rows: list[list[str]],
    *,
    numeric_columns: set[int] | None = None,
) -> None:
    """添加原生表格并按财务材料语义设置对齐。"""

    numeric_columns = numeric_columns or set()
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    widths = [w * ratio for ratio in (0.22, 0.42, 0.36)] if len(headers) == 3 else [w / len(headers)] * len(headers)
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for cidx, header in enumerate(headers):
        cell = table.cell(0, cidx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            apply_paragraph_style(paragraph, TextStyle(7.2, PAPER, "sans", bold=True, align=PP_ALIGN.CENTER))
    for ridx, row in enumerate(rows, start=1):
        for cidx, value in enumerate(row):
            cell = table.cell(ridx, cidx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ridx % 2 else TABLE_ALT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            align = PP_ALIGN.RIGHT if cidx in numeric_columns else PP_ALIGN.LEFT
            for paragraph in cell.text_frame.paragraphs:
                apply_paragraph_style(paragraph, TextStyle(7.3, INK, "sans", align=align, line_spacing=1.0))
    return shape


def slide_cover(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S01 封面。"""

    slide = blank_slide(prs)
    set_background(slide, INK)
    draw_fluid_background(slide, "dark")
    draw_dot_field(slide, "dark", 10.65, 1.1, 10, 18, 0.13)
    financials = data["financials"].set_index("fiscal_year")
    fy25 = financials.loc[2025]

    add_chrome(slide, "dark", 1, "AAPL / SEC 10-K / FY2021-FY2025", "EDITORIAL INK · NATIVE PPTX")
    add_text(slide, "25", 8.7, 0.75, 4.2, 2.6, TextStyle(132, blend(PAPER, INK, 0.12), "serif", bold=True, align=PP_ALIGN.RIGHT), margin=0)
    add_text(slide, "APPLE", 0.62, 1.16, 5.4, 0.55, TextStyle(11, blend(PAPER, INK, 0.65), "mono"), margin=0)
    add_text(slide, "历年财报\n解读", 0.54, 1.66, 5.9, 1.75, TextStyle(46, PAPER, "serif", bold=True, line_spacing=0.92), margin=0)
    add_text(slide, "2021–2025", 0.61, 3.42, 4.2, 0.40, TextStyle(22, COPPER, "serif", italic=True), margin=0)
    add_text(
        slide,
        "把 guizang HTML 的电子杂志风迁移为原生 PowerPoint 对象：大标题、细线、ghost number、数据大字报与可验证财务链路。",
        0.64,
        4.18,
        5.85,
        0.82,
        TextStyle(10.4, blend(PAPER, INK, 0.72), "sans", line_spacing=1.25),
        margin=0,
    )
    add_kpi_card(slide, 7.25, 3.04, 1.75, 1.18, "Revenue", f"{usd_bn(fy25['revenue']):.1f}B", "FY2025 net sales", "dark", COPPER)
    add_kpi_card(slide, 9.24, 3.04, 1.75, 1.18, "Net income", f"{usd_bn(fy25['net_income']):.1f}B", "profit recovery", "dark", CELADON)
    add_kpi_card(slide, 11.23, 3.04, 1.55, 1.18, "GM", f"{pct(fy25['gross_margin']):.1f}%", "gross margin", "dark", CORAL)
    add_year_axis(slide, 7.28, 5.05, 5.35, [2021, 2022, 2023, 2024, 2025], "dark")


def slide_thesis(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S02 核心判断页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER)
    draw_fluid_background(slide, "light")
    add_page_title(
        slide,
        "light",
        2,
        "EXECUTIVE MEMO",
        "FY2025 的恢复不是单点反弹，而是利润率、服务化与现金流共同支撑",
        "先看结论，再拆证据。页面风格保留杂志式留白与等宽元数据，数据全部来自处理后 CSV。",
    )
    financials = data["financials"].set_index("fiscal_year")
    regions = data["regions"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    fy24 = financials.loc[2024]
    fy25 = financials.loc[2025]
    cards = [
        ("收入回到新高", f"{usd_bn(fy25['revenue']):.1f}B", f"同比 {yoy(fy25['revenue'], fy24['revenue']) * 100:+.1f}%"),
        ("净利润恢复", f"{usd_bn(fy25['net_income']):.1f}B", f"同比 {yoy(fy25['net_income'], fy24['net_income']) * 100:+.1f}%"),
        ("毛利率上行", f"{pct(fy25['gross_margin']):.1f}%", "五年高点"),
        ("大中华区承压", f"{usd_mn_to_bn(regions.loc[2025, 'Greater China']):.1f}B", f"同比 {yoy(regions.loc[2025, 'Greater China'], regions.loc[2024, 'Greater China']) * 100:+.1f}%"),
    ]
    add_text(slide, "The machine lights up again.", 0.62, 2.28, 5.2, 0.46, TextStyle(21.5, INK, "serif", italic=True), margin=0)
    add_multiline_text(
        slide,
        [
            "FY2025 不是单一产品驱动的孤立恢复，而是收入、净利润、综合毛利率同时上修。",
            "Services 增长继续改善收入质量，现金流与资本回报仍然构成财务底座。",
            "风险侧没有消失，Greater China、AI 终端节奏、监管和硬件换机周期仍是后续跟踪重点。",
        ],
        0.67,
        2.96,
        5.25,
        1.34,
        TextStyle(9.5, INK_SOFT, "sans", line_spacing=1.33),
        bullet=False,
        margin=0,
    )
    add_statement_block(slide, 0.68, 4.84, 5.40, 0.86, "测试判断", "本页故意使用强叙事型 title 与少量正文，验证 native PPTX 能否承载 guizang 的演讲型 magazine rhythm。", "light")
    for idx, (label, value, note) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        accent = [COPPER, CELADON, INK, CORAL][idx]
        add_kpi_card(slide, 7.02 + col * 2.64, 2.30 + row * 1.42, 2.36, 1.14, label, value, note, "light", accent)
    add_year_axis(slide, 7.05, 5.68, 5.35, [2021, 2022, 2023, 2024, 2025], "light")


def slide_revenue_profit(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S03 收入与利润页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER_WARM)
    draw_dot_field(slide, "light", 0.70, 5.10, 26, 7, 0.11)
    add_page_title(
        slide,
        "light",
        3,
        "RESULT / TOP LINE",
        "收入创五年新高，净利润恢复到 1120 亿美元",
        "FY2025 收入同比增长 6.4%，净利润同比增长 19.5%，恢复强度明显高于 FY2024。",
    )
    financials = data["financials"].sort_values("fiscal_year")
    years = financials["fiscal_year"].astype(int).tolist()
    revenue = [usd_bn(v) for v in financials["revenue"]]
    net_income = [usd_bn(v) for v in financials["net_income"]]
    margin = [pct(v) for v in financials["net_margin"]]
    add_text(slide, "416.2B", 0.56, 2.42, 4.3, 0.80, TextStyle(43, INK, "serif", bold=True), margin=0)
    add_text(slide, "FY2025 Revenue", 0.70, 3.25, 2.4, 0.18, TextStyle(7.6, MUTED, "mono"), margin=0)
    add_text(slide, "+6.4% YoY", 3.45, 3.25, 1.05, 0.18, TextStyle(7.6, COPPER, "mono", align=PP_ALIGN.RIGHT, bold=True), margin=0)
    add_text(slide, "112.0B", 0.56, 3.80, 3.8, 0.58, TextStyle(31, INK_SOFT, "serif", bold=True), margin=0)
    add_text(slide, "FY2025 Net Income / +19.5% YoY", 0.70, 4.42, 3.4, 0.18, TextStyle(7.2, MUTED, "mono"), margin=0)
    add_native_bar_chart(slide, 5.08, 2.36, 3.05, 2.88, years, revenue, "light", color=COPPER, label="Revenue", value_format="{:.0f}")
    add_native_line_chart(
        slide,
        8.52,
        2.36,
        3.65,
        2.88,
        years,
        [("Net income", net_income, CELADON), ("Net margin", margin, CORAL)],
        "light",
        y_min=20,
        y_max=125,
        value_suffix="",
    )
    add_statement_block(slide, 5.10, 5.55, 7.05, 0.64, "读图结论", "2023 年收入回落后，2024-2025 连续恢复；FY2025 净利润增速高于收入增速，说明利润率与费用纪律共同放大了恢复。", "light")


def slide_product_mix(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S04 产品结构页。"""

    slide = blank_slide(prs)
    set_background(slide, INK)
    draw_fluid_background(slide, "dark")
    add_page_title(
        slide,
        "dark",
        4,
        "PRODUCT / SERVICE MIX",
        "iPhone 仍是收入中枢，Services 是最稳定的第二曲线",
        "FY2025 iPhone 收入约 2096 亿美元，Services 收入约 1092 亿美元并同比增长 13.5%。",
    )
    products = data["products"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    fy25 = products.loc[2025]
    parts = [
        ("iPhone", usd_mn_to_bn(fy25["iPhone"]), COPPER),
        ("Services", usd_mn_to_bn(fy25["Services"]), CELADON),
        ("Mac", usd_mn_to_bn(fy25["Mac"]), blend(PAPER, INK, 0.40)),
        ("iPad", usd_mn_to_bn(fy25["iPad"]), blend(PAPER, INK, 0.30)),
        ("Wearables", usd_mn_to_bn(fy25["Wearables, Home and Accessories"]), CORAL),
    ]
    add_stacked_strip(slide, 0.72, 2.55, 11.72, 0.78, parts, "dark")
    years = products.index.astype(int).tolist()
    iphone = [usd_mn_to_bn(v) for v in products["iPhone"]]
    services = [usd_mn_to_bn(v) for v in products["Services"]]
    add_native_line_chart(slide, 0.74, 4.25, 4.90, 1.55, years, [("iPhone", iphone, COPPER), ("Services", services, CELADON)], "dark", y_min=55, y_max=220, value_suffix="B")
    add_statement_block(slide, 6.18, 4.24, 2.72, 1.20, "结构含义", "Services 不是收入占比最高的业务，却是毛利率最强、增长最稳定的质量来源。", "dark")
    add_statement_block(slide, 9.28, 4.24, 2.72, 1.20, "风险含义", "iPhone 仍决定总收入中枢，硬件周期与区域需求会继续影响恢复斜率。", "dark")


def slide_margin_engine(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S05 毛利率引擎页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER)
    draw_fluid_background(slide, "light")
    add_page_title(
        slide,
        "light",
        5,
        "MARGIN ENGINE",
        "Services 的 75.4% 毛利率继续抬升整体质量",
        "Services 毛利率比 Products 高约 38.6pct，综合毛利率五年提升 5.1pct。",
    )
    gm = data["gross_margin"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    years = gm.index.astype(int).tolist()
    add_native_line_chart(
        slide,
        0.72,
        2.30,
        6.25,
        3.28,
        years,
        [
            ("Services", gm["Services"].tolist(), CELADON),
            ("Products", gm["Products"].tolist(), COPPER),
            ("Total", gm["Total gross margin percentage"].tolist(), CORAL),
        ],
        "light",
        y_min=30,
        y_max=80,
        value_suffix="%",
    )
    add_text(slide, "75.4%", 7.52, 2.36, 2.10, 0.56, TextStyle(34, INK, "serif", bold=True), margin=0)
    add_text(slide, "Services gross margin", 7.60, 2.95, 2.3, 0.16, TextStyle(6.8, MUTED, "mono"), margin=0)
    add_text(slide, "36.8%", 10.12, 2.36, 1.62, 0.56, TextStyle(27, INK_SOFT, "serif", bold=True), margin=0)
    add_text(slide, "Products gross margin", 10.18, 2.95, 2.1, 0.16, TextStyle(6.8, MUTED, "mono"), margin=0)
    add_line(slide, 9.72, 2.42, 9.72, 3.30, blend(INK, PAPER, 0.20), 0.55)
    add_statement_block(slide, 7.62, 3.76, 4.30, 0.82, "机制解释", "当 Services 占比继续提升，综合毛利率获得结构性支撑；这不是价格单点，而是生态、订阅和平台分成共同形成的质量改善。", "light")
    add_statement_block(slide, 7.62, 4.86, 4.30, 0.70, "费用侧", "R&D ratio 从 FY2021 的 6.0% 提升到 FY2025 的 8.3%，说明 Apple 仍在用研发强度换下一轮产品周期。", "light")


def slide_regions(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S06 地区结构页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER_WARM)
    add_page_title(
        slide,
        "light",
        6,
        "GEOGRAPHY / MIX",
        "美欧日及亚太扩张，大中华区仍是结构性压力点",
        "FY2025 Greater China 收入同比下降 3.8%，其他主要地区均为正增长。",
    )
    regions = data["regions"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    rows = []
    for name in ["Americas", "Europe", "Greater China", "Japan", "Rest of Asia Pacific"]:
        value = usd_mn_to_bn(regions.loc[2025, name])
        growth = yoy(regions.loc[2025, name], regions.loc[2024, name]) * 100
        rows.append((name, value, growth))
    rows.sort(key=lambda item: item[1], reverse=True)
    add_horizontal_bar_matrix(slide, 0.72, 2.32, 6.62, 3.32, rows, "light")
    add_text(slide, "Greater China", 8.08, 2.42, 3.7, 0.38, TextStyle(22, CORAL, "serif", bold=True), margin=0)
    add_text(slide, "不是规模不重要，而是 FY2023-FY2025 连续没有恢复到 FY2022 高点。", 8.16, 2.90, 3.95, 0.46, TextStyle(9.4, INK_SOFT, "sans", line_spacing=1.25), margin=0)
    add_statement_block(slide, 8.18, 3.86, 3.82, 0.72, "正向支撑", "Americas 与 Europe 是 FY2025 增量主轴，日本和 Rest of Asia Pacific 也同步回升。", "light")
    add_statement_block(slide, 8.18, 4.88, 3.82, 0.72, "风险变量", "竞争、渠道、宏观消费和监管都可能放大区域波动，需单独跟踪。", "light")


def slide_cash_machine(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S07 现金机器页。"""

    slide = blank_slide(prs)
    set_background(slide, INK)
    draw_fluid_background(slide, "dark")
    add_page_title(
        slide,
        "dark",
        7,
        "CASH MACHINE",
        "自由现金流仍厚，资本回报强度没有明显退场",
        "FY2025 自由现金流约 987.7 亿美元，回购与分红合计约 1061.3 亿美元。",
    )
    financials = data["financials"].sort_values("fiscal_year")
    years = financials["fiscal_year"].astype(int).tolist()
    fcf = [usd_bn(v) for v in financials["fcf"]]
    capital_return = [usd_bn(v) for v in financials["capital_return"]]
    net_cash = [usd_bn(v) for v in financials["net_cash"]]
    fy25 = financials[financials["fiscal_year"] == 2025].iloc[0]
    add_text(slide, f"{usd_bn(fy25['fcf']):.1f}B", 0.64, 2.32, 3.70, 0.78, TextStyle(39, PAPER, "serif", bold=True), margin=0)
    add_text(slide, "FY2025 Free Cash Flow", 0.75, 3.16, 2.7, 0.16, TextStyle(7.0, blend(PAPER, INK, 0.60), "mono"), margin=0)
    add_text(slide, f"{usd_bn(fy25['capital_return']):.1f}B", 0.64, 3.76, 3.70, 0.58, TextStyle(30, COPPER, "serif", bold=True), margin=0)
    add_text(slide, "Buybacks + Dividends", 0.75, 4.38, 2.7, 0.16, TextStyle(7.0, blend(PAPER, INK, 0.60), "mono"), margin=0)
    add_native_bar_chart(slide, 4.82, 2.42, 3.24, 2.62, years, fcf, "dark", color=CELADON, label="Free cash flow", value_format="{:.0f}")
    add_native_line_chart(slide, 8.42, 2.42, 3.62, 2.62, years, [("Capital return", capital_return, COPPER), ("Net cash", net_cash, CORAL)], "dark", y_min=20, y_max=115, value_suffix="B")
    add_statement_block(slide, 4.85, 5.42, 7.18, 0.62, "资本配置口径", "资本回报 = buybacks + dividends；净现金 = cash & marketable securities - total debt。FY2025 净现金仍为正，但较 FY2021 明显收敛。", "dark")


def slide_pipeline(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S08 数据管线页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER)
    draw_dot_field(slide, "light", 9.86, 1.18, 16, 24, 0.10)
    add_page_title(
        slide,
        "light",
        8,
        "DATA PIPELINE",
        "从 SEC 原始披露到页面判断，数据链路必须可复跑",
        "本页只参考 apple-financial-report-review 的数据管线，不参考旧 PPT 的具体视觉表现。",
    )
    stages = [
        ("01", "SEC APIs", "companyfacts / submissions"),
        ("02", "10-K HTML", "FY2021-FY2025 annual reports"),
        ("03", "Extract", "segment tables + xbrl fields"),
        ("04", "Processed CSV", "financials / products / regions / margins"),
        ("05", "Narrative", "brief + slide contracts + asset slots"),
        ("06", "Native PPTX", "editable shapes + tables + validation"),
    ]
    start_x, y = 0.72, 2.58
    card_w, card_h, gap = 1.78, 1.02, 0.22
    for idx, (num, title, desc) in enumerate(stages):
        x = start_x + idx * (card_w + gap)
        fill = INK if idx in {0, 5} else WHITE
        text = PAPER if idx in {0, 5} else INK
        accent = COPPER if idx < 3 else CELADON
        add_rect(slide, x, y, card_w, card_h, fill, line=blend(INK, PAPER, 0.18), line_width=0.5)
        add_text(slide, num, x + 0.10, y + 0.10, 0.42, 0.18, TextStyle(8.0, accent, "mono", bold=True), margin=0)
        add_text(slide, title, x + 0.10, y + 0.36, card_w - 0.20, 0.22, TextStyle(10.0, text, "serif", bold=True), margin=0)
        add_text(slide, desc, x + 0.10, y + 0.68, card_w - 0.20, 0.22, TextStyle(6.4, blend(text, fill, 0.70), "sans", line_spacing=1.1), margin=0)
        if idx < len(stages) - 1:
            add_line(slide, x + card_w + 0.04, y + card_h / 2, x + card_w + gap - 0.04, y + card_h / 2, COPPER, 0.8)
    add_statement_block(slide, 0.78, 4.45, 5.35, 0.82, "正确地失败", "数据管线只复用本地已有 CSV 与来源记录。若数据文件缺失、字段缺失或导出失败，脚本应直接暴露错误，不用占位数字伪装成功。", "light")
    add_statement_block(slide, 6.70, 4.45, 5.35, 0.82, "风格迁移", "HTML 中的 WebGL 背景被替换为 PowerPoint 原生椭圆、线条与点阵；页面内容仍保持可编辑和可结构检查。", "light")


def slide_watchlist(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S09 后续观察变量页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER_WARM)
    draw_fluid_background(slide, "light")
    add_chrome(slide, "light", 9, "NEXT VARIABLES / AFTER FY2025", "WATCHLIST · 09 / 10")
    add_text(slide, "NEXT", 0.55, 0.95, 5.20, 0.94, TextStyle(54, INK, "serif", bold=True), margin=0)
    add_text(
        slide,
        "财务质量已经修复，再加速要看 AI 终端、服务生态与区域竞争",
        0.65,
        2.02,
        7.2,
        0.44,
        TextStyle(18, INK_SOFT, "serif", bold=True, line_spacing=1.0),
        margin=0,
    )
    cards = [
        ("01", "AI 终端换机", "是否能把 AI 功能转化为硬件换机理由，而不是只停留在软件层体验。", COPPER),
        ("02", "Services monetization", "高毛利 Services 能否继续增长，同时处理 App Store、订阅和监管压力。", CELADON),
        ("03", "区域竞争与监管", "Greater China、欧盟监管、供应链与汇率共同决定增长质量和风险折价。", CORAL),
    ]
    for idx, (num, title, body, accent) in enumerate(cards):
        x = 0.78 + idx * 4.03
        y = 3.20
        add_rect(slide, x, y, 3.48, 1.70, WHITE, line=blend(INK, PAPER, 0.18), line_width=0.5)
        add_text(slide, num, x + 0.16, y + 0.18, 0.54, 0.28, TextStyle(13, accent, "mono", bold=True), margin=0)
        add_text(slide, title, x + 0.16, y + 0.62, 2.9, 0.24, TextStyle(13, INK, "serif", bold=True), margin=0)
        add_text(slide, body, x + 0.16, y + 1.04, 2.95, 0.45, TextStyle(8.4, INK_SOFT, "sans", line_spacing=1.18), margin=0)
    add_text(slide, "Research frame, not rating.", 7.88, 1.22, 4.8, 0.40, TextStyle(22, COPPER, "serif", italic=True, align=PP_ALIGN.RIGHT), margin=0)


def slide_sources_risks(prs: Presentation, data: dict[str, pd.DataFrame | dict]) -> None:
    """生成 S10 来源与风险页。"""

    slide = blank_slide(prs)
    set_background(slide, PAPER)
    add_page_title(
        slide,
        "light",
        10,
        "SOURCES / RISK",
        "来源、口径与风险声明",
        "本材料仅供学术交流和 skill 测试使用，不构成投资建议。",
    )
    sources = data["sources"].sort_values("fiscal_year", ascending=False)
    source_rows = [
        [f"FY{int(row.fiscal_year)}", str(row.filing_date), str(row.accession)]
        for row in sources.itertuples()
    ]
    add_text(slide, "已下载 10-K", 0.72, 2.28, 2.0, 0.18, TextStyle(9.2, INK, "serif", bold=True), margin=0)
    add_native_table(slide, 0.72, 2.58, 5.18, 1.70, ["FY", "Filed", "Accession"], source_rows, numeric_columns=set())
    risks = [
        ["经营", "产品创新节奏、硬件换机周期、Services 增速与生态抽成变化"],
        ["区域", "Greater China 需求、竞争格局、渠道库存与宏观消费"],
        ["监管", "App Store、数据隐私、反垄断、税务与跨境业务约束"],
        ["财务", "汇率、资本回报节奏、利率、供应链成本和市场估值波动"],
    ]
    add_text(slide, "风险矩阵", 6.55, 2.28, 2.0, 0.18, TextStyle(9.2, INK, "serif", bold=True), margin=0)
    add_native_table(slide, 6.55, 2.58, 5.70, 1.70, ["类型", "说明", "口径"], [[r[0], r[1], "需持续跟踪"] for r in risks], numeric_columns=set())
    add_statement_block(slide, 0.72, 4.78, 5.20, 0.82, "数据边界", "核心三表数据来自 SEC companyfacts；产品线、地区与毛利率来自 Apple 10-K HTML 表格抽取。金额默认单位为十亿美元。", "light")
    add_statement_block(slide, 6.55, 4.78, 5.70, 0.82, "免责声明", "本 deck 是原生 PPTX 构建与视觉风格迁移测试，不代表任何组织机构，机构和个人的观点和立场，不构成投资建议、评级或目标价。", "light")


def build_deck() -> Presentation:
    """构建完整 deck。"""

    data = load_inputs()
    prs = make_prs()
    builders = [
        slide_cover,
        slide_thesis,
        slide_revenue_profit,
        slide_product_mix,
        slide_margin_engine,
        slide_regions,
        slide_cash_machine,
        slide_pipeline,
        slide_watchlist,
        slide_sources_risks,
    ]
    for builder in builders:
        builder(prs, data)
    return prs


def write_summary(prs: Presentation) -> None:
    """写出构建摘要。"""

    SUMMARY_PATH.parent.mkdir(parents=True, exist_ok=True)
    summary = {
        "generated_at": date.today().isoformat(),
        "workspace": str(WORKSPACE),
        "slide_count": len(prs.slides),
        "pptx": str(OUT_PPTX),
        "final_pptx": str(FINAL_PPTX),
        "style_profile": "guizang Style A editorial_ink, indigo porcelain native PowerPoint translation",
        "data_files": sorted(path.name for path in DATA_DIR.glob("*.csv")),
    }
    SUMMARY_PATH.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")


def fix_docprops_slide_count(pptx_path: Path, slide_count: int) -> None:
    """修正 `docProps/app.xml` 中的 Slides 统计，避免 package preflight 误报。"""

    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        with zipfile.ZipFile(pptx_path, "r") as source_zip:
            source_zip.extractall(temp_path)
        app_xml = temp_path / "docProps" / "app.xml"
        tree = ET.parse(app_xml)
        root = tree.getroot()
        slides = root.find("ep:Slides", ns)
        if slides is None:
            slides = ET.SubElement(root, f"{{{ns['ep']}}}Slides")
        slides.text = str(slide_count)
        ET.register_namespace("", ns["ep"])
        tree.write(app_xml, encoding="UTF-8", xml_declaration=True)

        fixed_path = pptx_path.with_suffix(".fixed.pptx")
        with zipfile.ZipFile(fixed_path, "w", compression=zipfile.ZIP_DEFLATED) as target_zip:
            for file_path in temp_path.rglob("*"):
                if file_path.is_file():
                    target_zip.write(file_path, file_path.relative_to(temp_path).as_posix())
        shutil.move(fixed_path, pptx_path)


def main() -> int:
    """执行 PPTX 构建。"""

    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT_PPTX)
    fix_docprops_slide_count(OUT_PPTX, len(prs.slides))
    shutil.copy2(OUT_PPTX, FINAL_PPTX)
    write_summary(prs)
    print(f"[OK] 写出 PPTX: {OUT_PPTX}")
    print(f"[OK] 写出 final: {FINAL_PPTX}")
    print(f"[INFO] slides={len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
